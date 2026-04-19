"""
VibeCoding Workshop 投影片生成器
簡約專業風格 — 淺色背景、微軟正黑體
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── 色彩系統（淺色主題）─────────────────────────────────
BG_DARK      = RGBColor(0xFA, 0xFA, 0xFC)  # 近白背景
BG_CODE      = RGBColor(0xF0, 0xF2, 0xF5)  # 程式碼區塊淺灰底
ACCENT       = RGBColor(0x1A, 0x6B, 0xE5)  # 藍色強調
ACCENT2      = RGBColor(0x1B, 0x8C, 0x4E)  # 綠色（通過/成功）
ACCENT_WARN  = RGBColor(0xC4, 0x8A, 0x0A)  # 深黃（警告）
ACCENT_ERR   = RGBColor(0xD1, 0x2E, 0x2E)  # 紅色（失敗）
WHITE        = RGBColor(0x1A, 0x1A, 0x1A)  # 主文字（深黑）
GRAY_LIGHT   = RGBColor(0x37, 0x37, 0x3C)  # 次要文字
GRAY_MID     = RGBColor(0x6B, 0x6B, 0x76)  # 輔助文字
GRAY_DIM     = RGBColor(0xA0, 0xA0, 0xAA)  # 最淡文字
PURE_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)  # 純白（裝飾用）

# ── 字體 ─────────────────────────────────────────────
FONT_TITLE   = "Microsoft JhengHei UI"  # 微軟正黑體
FONT_BODY    = "Microsoft JhengHei UI"
FONT_CODE    = "Microsoft JhengHei UI"
FONT_CJK     = "Microsoft JhengHei UI"

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_font(run, name=FONT_BODY, size=Pt(18), color=WHITE, bold=False, italic=False):
    f = run.font
    f.name = name
    f.east_asian_font = FONT_CJK
    f.size = size
    f.color.rgb = color
    f.bold = bold
    f.italic = italic


def add_paragraph(tf, text, font_name=FONT_BODY, size=Pt(18), color=WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=size, color=color, bold=bold)
    return p


def make_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # 標題
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_font(run, name=FONT_TITLE, size=Pt(44), color=WHITE, bold=True)

    if subtitle:
        add_paragraph(tf, subtitle, size=Pt(22), color=GRAY_MID,
                      alignment=PP_ALIGN.CENTER, space_before=Pt(16))

    # 底部裝飾線
    from pptx.shapes.autoshape import Shape
    line = slide.shapes.add_shape(
        1, Inches(4), Inches(3.7), Inches(5.3), Pt(3)  # 1 = rectangle
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return slide


def make_section_slide(prs, chapter, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = chapter
    set_font(run, size=Pt(20), color=ACCENT)

    tb2 = add_text_box(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = title
    set_font(run2, name=FONT_TITLE, size=Pt(40), color=WHITE, bold=True)
    return slide


def make_content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 標題
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, name=FONT_TITLE, size=Pt(32), color=WHITE, bold=True)

    # 底線
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.25), Inches(11.7), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # 內容
    tb2 = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    for item in bullets:
        if isinstance(item, tuple):
            text, color, sz, bld = item
        else:
            text, color, sz, bld = item, WHITE, Pt(20), False
        add_paragraph(tf2, text, size=sz, color=color, bold=bld, space_after=Pt(8))

    # 講師備註
    if note:
        from pptx.util import Emu
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = note

    return slide


def make_code_slide(prs, title, code_lines, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 標題
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, name=FONT_TITLE, size=Pt(30), color=WHITE, bold=True)

    # 程式碼區塊背景
    code_bg = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.7)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BG_CODE
    code_bg.line.color.rgb = RGBColor(0xD8, 0xDA, 0xE0)
    code_bg.line.width = Pt(1)
    # 圓角效果（近似）
    code_bg.shadow.inherit = False

    # 程式碼文字
    tb2 = add_text_box(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.3))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    for i, line in enumerate(code_lines):
        if isinstance(line, tuple):
            text, color = line
        else:
            text, color = line, GRAY_LIGHT

        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = text
        set_font(run, name=FONT_CODE, size=Pt(16), color=color)

    if note:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = note

    return slide


def make_two_col_slide(prs, title, left_title, left_items, right_title, right_items, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 標題
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, name=FONT_TITLE, size=Pt(32), color=WHITE, bold=True)

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.25), Inches(11.7), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # 左欄
    tb_l = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    run_l = p_l.add_run()
    run_l.text = left_title
    set_font(run_l, size=Pt(22), color=ACCENT, bold=True)
    for item in left_items:
        if isinstance(item, tuple):
            text, color = item
        else:
            text, color = item, WHITE
        add_paragraph(tf_l, text, size=Pt(18), color=color, space_after=Pt(6))

    # 右欄
    tb_r = add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    run_r = p_r.add_run()
    run_r.text = right_title
    set_font(run_r, size=Pt(22), color=ACCENT2, bold=True)
    for item in right_items:
        if isinstance(item, tuple):
            text, color = item
        else:
            text, color = item, WHITE
        add_paragraph(tf_r, text, size=Pt(18), color=color, space_after=Pt(6))

    if note:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = note

    return slide


def make_checklist_slide(prs, title, items, note=""):
    """items: list of (text, checked:bool)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, name=FONT_TITLE, size=Pt(32), color=WHITE, bold=True)

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.25), Inches(11.7), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    tb2 = add_text_box(slide, Inches(1.2), Inches(1.6), Inches(11.0), Inches(5.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    for text, checked in items:
        icon = "\u2705" if checked else "\u2610"
        color = ACCENT2 if checked else GRAY_LIGHT
        add_paragraph(tf2, f"  {icon}  {text}", size=Pt(22), color=color, space_after=Pt(12))

    if note:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = note

    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ================================================================
    # COVER
    # ================================================================
    make_title_slide(prs,
        "VibeCoding 模板實戰 Workshop",
        "用 Claude Code 跑完一個專案的開發工作流"
    )

    # ================================================================
    # CHAPTER 1 — 認識模板
    # ================================================================
    make_section_slide(prs, "Chapter 1", "認識 VibeCoding 模板")

    # 痛點
    make_content_slide(prs, "用 AI 寫程式的 5 個痛點", [
        ("1.  每次對話都要重新解釋專案背景", ACCENT_ERR, Pt(22), False),
        ("2.  AI 產出的程式碼品質不穩定", ACCENT_ERR, Pt(22), False),
        ("3.  沒有測試、沒有審查，直接上線", ACCENT_ERR, Pt(22), False),
        ("4.  需求 \u2192 設計 \u2192 實作之間斷裂", ACCENT_ERR, Pt(22), False),
        ("5.  換個專案又從零開始設定", ACCENT_ERR, Pt(22), False),
        ("", WHITE, Pt(12), False),
        ("\u2192  問題不在 AI 不夠聰明，而是缺乏系統化工作流", GRAY_MID, Pt(18), False),
    ], note="互動：請學員舉手看看遇過幾個")

    # 模板是什麼
    make_content_slide(prs, "VibeCoding 模板 = AI 開發的標準作業程序", [
        ("\U0001f4cb  16 份文件模板        規範從需求到部署的每個階段", WHITE, Pt(20), False),
        ("\U0001f916  13 個 AI Agent        各司其職、自動協作", WHITE, Pt(20), False),
        ("\u26a1  17 個快捷指令        一行啟動完整流程", WHITE, Pt(20), False),
        ("\U0001f4cf   9 條自動規則         品質底線自動守護", WHITE, Pt(20), False),
        ("\U0001f527   7 個領域技能         按需載入專業知識", WHITE, Pt(20), False),
        ("", WHITE, Pt(12), False),
        ("\u2192  複製到任何新專案，立即可用", ACCENT, Pt(20), True),
    ])

    # 對比
    make_two_col_slide(prs, "對比 — 沒模板 vs 有模板",
        "\u274c  沒有模板", [
            "\"幫我寫一個 Todo App\"",
            "\u2192 AI 直接寫一大段程式碼",
            "\u2192 沒測試",
            "\u2192 沒文件",
            "\u2192 沒品質把關",
            "\u2192 出問題才發現",
        ],
        "\u2705  有模板", [
            "/task-init  \u2192 互動式 Q&A",
            "/task-next  \u2192 拆好的任務",
            "/plan       \u2192 先規劃再動手",
            "/tdd        \u2192 先寫測試再實作",
            "/review-code \u2192 AI 審查",
            "/verify     \u2192 品質門把關",
        ],
        note="核心訊息：差別在流程，不在 AI 模型本身"
    )

    # 5 大機制
    make_content_slide(prs, "5 大擴展機制", [
        ("Hook        系統事件自動觸發 \u2192 Shell 腳本", WHITE, Pt(20), False),
        ("                  例：啟動時偵測模板、自動追蹤時間", GRAY_MID, Pt(16), False),
        ("Command   使用者輸入 /xxx \u2192 預設 Prompt", WHITE, Pt(20), False),
        ("                  例：/tdd, /plan, /verify", GRAY_MID, Pt(16), False),
        ("Skill          語意偵測自動觸發 \u2192 領域知識", WHITE, Pt(20), False),
        ("                  例：偵測到 FastAPI 自動載入相關知識", GRAY_MID, Pt(16), False),
        ("Agent        主 Agent 委派子任務 \u2192 專業 AI", WHITE, Pt(20), False),
        ("                  例：code-quality-specialist 做審查", GRAY_MID, Pt(16), False),
        ("Context     跨 Agent 知識共享 \u2192 報告存檔", WHITE, Pt(20), False),
        ("                  例：品質報告 \u2192 測試報告接力", GRAY_MID, Pt(16), False),
    ], note="今天主要用 Command 和 Agent，其他三個先知道就好")

    # 13 Agent
    make_content_slide(prs, "13 個 Agent — 按模型分級", [
        ("Opus  (\u91cd\u91cf\u7d1a\u63a8\u7406)", ACCENT, Pt(22), True),
        ("    planner \u2022 architect \u2022 security-auditor", WHITE, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("Sonnet  (\u4e00\u822c\u958b\u767c)", ACCENT2, Pt(22), True),
        ("    code-quality \u2022 tdd-guide \u2022 e2e-specialist", WHITE, Pt(18), False),
        ("    test-engineer \u2022 refactor-cleaner \u2022 deployment", WHITE, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("Haiku  (\u8f15\u91cf\u5feb\u901f)", ACCENT_WARN, Pt(22), True),
        ("    build-resolver \u2022 doc-specialist \u2022 template-mgr", WHITE, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("\u2192  \u6a21\u677f\u81ea\u52d5\u5206\u914d\uff0c\u4e0d\u9700\u8981\u624b\u52d5\u9078\u64c7", GRAY_MID, Pt(16), False),
    ])

    # 目錄結構
    make_code_slide(prs, "目錄結構 — .claude/ 下的關鍵目錄", [
        (".claude/", ACCENT),
        ("\u251c\u2500\u2500 agents/          \u2190 13 \u500b Agent \u5b9a\u7fa9", GRAY_LIGHT),
        ("\u251c\u2500\u2500 commands/        \u2190 17 \u500b\u5feb\u6377\u6307\u4ee4", GRAY_LIGHT),
        ("\u251c\u2500\u2500 rules/           \u2190 9 \u689d\u81ea\u52d5\u898f\u5247\uff08\u6bcf\u6b21\u5c0d\u8a71\u8f09\u5165\uff09", ACCENT2),
        ("\u251c\u2500\u2500 skills/          \u2190 7 \u500b\u9818\u57df\u6280\u80fd", GRAY_LIGHT),
        ("\u251c\u2500\u2500 hooks/           \u2190 6 \u500b\u7cfb\u7d71\u4e8b\u4ef6\u8173\u672c", GRAY_LIGHT),
        ("\u251c\u2500\u2500 context/         \u2190 Agent \u9593\u7684\u5831\u544a\u5171\u4eab", GRAY_LIGHT),
        ("\u251c\u2500\u2500 guides/          \u2190 \u53c3\u8003\u6587\u4ef6\uff08\u4e0d\u81ea\u52d5\u8f09\u5165\uff09", GRAY_DIM),
        ("\u2514\u2500\u2500 taskmaster-data/ \u2190 WBS \u4efb\u52d9 + \u6642\u9593\u8ffd\u8e64", ACCENT),
    ], note="不需要背，知道去哪找就好")

    # 環境準備
    make_checklist_slide(prs, "環境準備檢查", [
        ("Claude Code CLI 已安裝    \u2192 claude --version", False),
        ("Python 3.11+ 已安裝         \u2192 python --version", False),
        ("Node.js 18+ 已安裝           \u2192 node --version", False),
        ("Git 已設定                          \u2192 git config user.name", False),
        ("模板已複製到工作目錄", False),
        (".mcp.json 已設定 API Key", False),
    ], note="給學員 1 分鐘檢查，有問題的舉手")

    # 旅程
    make_content_slide(prs, "今天的旅程", [
        ("Ch1   認識模板                10 min  \u2190 你在這裡", ACCENT, Pt(22), True),
        ("Ch2   專案初始化             15 min", WHITE, Pt(22), False),
        ("Ch3   開發循環 x3            25 min   \u2b50 核心", ACCENT_WARN, Pt(22), True),
        ("Ch4   品質驗證                10 min", WHITE, Pt(22), False),
        ("Ch5   進階 + Q&A             10 min", WHITE, Pt(22), False),
        ("", WHITE, Pt(12), False),
        ("\u2192  70 分鐘後，你會了一套 AI 開發工作流", ACCENT2, Pt(20), True),
    ])

    # ================================================================
    # CHAPTER 2 — 專案初始化
    # ================================================================
    make_section_slide(prs, "Chapter 2", "專案初始化")

    # /task-init
    make_content_slide(prs, "一切從 /task-init 開始", [
        ("你輸入：  /task-init TaskFlow", ACCENT, Pt(24), True),
        ("", WHITE, Pt(8), False),
        ("AI 會做：", WHITE, Pt(20), True),
        ("  1.  問你一系列問題（互動式 Q&A）", WHITE, Pt(20), False),
        ("        \u2022 專案做什麼？用什麼技術？範圍多大？", GRAY_MID, Pt(16), False),
        ("  2.  根據回答自動產出：", WHITE, Pt(20), False),
        ("        \u2022 CLAUDE.md  \u2014 專案說明書", ACCENT2, Pt(18), False),
        ("        \u2022 wbs.md     \u2014 任務拆解", ACCENT2, Pt(18), False),
        ("  3.  從此 AI 每次對話都知道你的專案脈絡", WHITE, Pt(20), False),
    ], note="這是最關鍵的 Aha Moment — AI 先問問題、理解需求，再產文件")

    # MVP vs Complete
    make_two_col_slide(prs, "MVP vs Complete 模式",
        "MVP 模式  \u2190 今天用這個", [
            "輕量 Tech Spec",
            "快速驗證想法",
            "1-2 人小專案",
            "原型 / POC",
            "",
            ("隨時可升級為 Complete", GRAY_MID),
        ],
        "Complete 模式", [
            "完整 16 份文件",
            "每個決策有文件記錄",
            "3+ 人團隊專案",
            "正式產品 / 合規需求",
            "",
            ("大專案、正式產品用這個", GRAY_MID),
        ]
    )

    # CLAUDE.md
    make_content_slide(prs, "CLAUDE.md = 專案的「身分證」", [
        ("AI 每次對話都會自動讀取這份檔案", ACCENT, Pt(20), True),
        ("", WHITE, Pt(8), False),
        ("\u2022  知道專案用什麼技術棧", WHITE, Pt(20), False),
        ("\u2022  知道程式碼規範（命名、風格）", WHITE, Pt(20), False),
        ("\u2022  知道目錄結構", WHITE, Pt(20), False),
        ("\u2022  知道測試策略", WHITE, Pt(20), False),
        ("\u2022  知道部署方式", WHITE, Pt(20), False),
        ("", WHITE, Pt(12), False),
        ("\u2192  你不需要每次重新解釋", ACCENT2, Pt(20), False),
        ("\u2192  團隊成員接手也能立刻上手", ACCENT2, Pt(20), False),
    ])

    # WBS
    make_code_slide(prs, "WBS = 自動拆解的任務清單", [
        ("1.0 \u5c08\u6848\u7ba1\u7406", ACCENT),
        ("\u251c\u2500\u2500 1.1 \u9700\u6c42\u5206\u6790          \u2190 /task-init \u7522\u51fa", GRAY_LIGHT),
        ("\u2514\u2500\u2500 1.2 \u67b6\u69cb\u8a2d\u8a08", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("2.0 \u5f8c\u7aef\u958b\u767c", ACCENT),
        ("\u251c\u2500\u2500 2.1 \u57fa\u790e\u67b6\u69cb          \u2190 /task-next \u6703\u6307\u6d3e", ACCENT2),
        ("\u251c\u2500\u2500 2.2 CRUD API", GRAY_LIGHT),
        ("\u2514\u2500\u2500 2.3 \u7be9\u9078\u529f\u80fd", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("3.0 \u524d\u7aef\u958b\u767c", ACCENT),
        ("\u251c\u2500\u2500 3.1 \u5143\u4ef6\u958b\u767c", GRAY_LIGHT),
        ("\u2514\u2500\u2500 3.2 \u4e92\u52d5\u529f\u80fd", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("\u6bcf\u500b\u4efb\u52d9\u6709\uff1a\u7de8\u865f \u2502 \u8ca0\u8cac\u4eba \u2502 \u72c0\u614b \u2502 \u5de5\u6642", GRAY_MID),
    ])

    # 操作指引
    make_code_slide(prs, "\u26a1 \u73fe\u5834\u64cd\u4f5c \u2014 /task-init", [
        ("# 1. \u6253\u958b\u7d42\u7aef\u6a5f", GRAY_MID),
        ("cd /your/workspace", WHITE),
        ("claude", ACCENT),
        ("", GRAY_DIM),
        ("# 2. \u57f7\u884c\u521d\u59cb\u5316", GRAY_MID),
        ("/task-init TaskFlow", ACCENT),
        ("", GRAY_DIM),
        ("# 3. \u56de\u7b54 AI \u7684\u554f\u984c", GRAY_MID),
        ("\u5c08\u6848\u63cf\u8ff0\uff1a\u6975\u7c21\u5f85\u8fa6\u6e05\u55ae\u61c9\u7528", WHITE),
        ("\u6280\u8853\u68e7\uff1a  Python FastAPI + React + SQLite", WHITE),
        ("\u6a21\u5f0f\uff1a    MVP", WHITE),
        ("\u529f\u80fd\uff1a    CRUD + \u5206\u985e\u7be9\u9078", WHITE),
        ("", GRAY_DIM),
        ("# 4. \u6aa2\u67e5\u7522\u51fa", GRAY_MID),
        ("/task-status", ACCENT),
    ], note="讓學員跟著做。注意 AI 一次只問一題。")

    # ================================================================
    # CHAPTER 3 — 開發循環
    # ================================================================
    make_section_slide(prs, "Chapter 3", "開發循環")

    # 循環全貌
    make_content_slide(prs, "開發循環全貌", [
        ("/task-next      取得下一個任務", WHITE, Pt(22), False),
        ("      \u2193", GRAY_DIM, Pt(16), False),
        ("/plan               planner agent (Opus) 建立計畫", WHITE, Pt(22), False),
        ("      \u2193            \u2192 人類審核 \u2192 確認", ACCENT_WARN, Pt(16), False),
        ("/tdd                tdd-guide agent \u2014 RED \u2192 GREEN \u2192 IMPROVE", WHITE, Pt(22), False),
        ("      \u2193", GRAY_DIM, Pt(16), False),
        ("/review-code   code-quality agent 審查", WHITE, Pt(22), False),
        ("      \u2193            \u2192 修復 CRITICAL / HIGH", ACCENT_WARN, Pt(16), False),
        ("git commit        conventional commits 格式", WHITE, Pt(22), False),
        ("      \u2193", GRAY_DIM, Pt(16), False),
        ("\u2500\u2500\u2500 \u9084\u6709\u4efb\u52d9\uff1f\u2192 \u56de\u5230 /task-next \u2500\u2500\u2500", ACCENT, Pt(18), True),
    ], note="核心觀念：人類主導、AI 輔助。/plan 後一定要人類確認。")

    # 第 1 輪
    make_section_slide(prs, "\u7b2c 1 \u8f2a\uff08\u8a73\u7d30\uff09", "\u5f8c\u7aef API \u2014 FastAPI + pytest")

    # /plan
    make_code_slide(prs, "/plan \u2014 AI \u5148\u60f3\u6e05\u695a\u518d\u52d5\u624b", [
        ("\U0001f4cb \u5be6\u4f5c\u8a08\u756b                            planner agent (Opus)", ACCENT),
        ("", GRAY_DIM),
        ("\u6b65\u9a5f 1: \u5efa\u7acb\u5c08\u6848\u7d50\u69cb", WHITE),
        ("  - pyproject.toml (dependencies)", GRAY_LIGHT),
        ("  - app/main.py    (FastAPI entry)", GRAY_LIGHT),
        ("  - app/models.py  (SQLModel schemas)", GRAY_LIGHT),
        ("  - app/database.py(SQLite connection)", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("\u6b65\u9a5f 2: \u5b9a\u7fa9 Todo \u8cc7\u6599\u6a21\u578b", WHITE),
        ("  - id, title, category, done, created_at", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("\u6b65\u9a5f 3: \u5be6\u4f5c CRUD endpoints", WHITE),
        ("  - GET / POST / PUT / DELETE  /api/todos", GRAY_LIGHT),
        ("", GRAY_DIM),
        ("\u98a8\u96aa: \u7121\uff08\u6a19\u6e96 CRUD\uff0c\u4f4e\u8907\u96dc\u5ea6\uff09", ACCENT2),
    ], note="暫停讓學員讀計畫。強調：AI 不是直接寫 code，而是先用 Opus 做規劃。")

    # /tdd
    make_content_slide(prs, "/tdd \u2014 \u5148\u5beb\u6e2c\u8a66\uff0c\u518d\u5beb\u5be6\u4f5c", [
        ("tdd-guide agent (Sonnet) \u5f37\u5236\u57f7\u884c\uff1a", GRAY_MID, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("  RED          \u5beb\u6e2c\u8a66\uff0c\u57f7\u884c \u2192 \u5168\u90e8 FAIL", ACCENT_ERR, Pt(24), True),
        ("                     pytest \u2192 \u274c \u7d05\u8272\uff08\u9019\u662f\u5c0d\u7684\uff01\uff09", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(8), False),
        ("  GREEN      \u5beb\u6700\u5c0f\u5be6\u4f5c\uff0c\u57f7\u884c \u2192 \u5168\u90e8 PASS", ACCENT2, Pt(24), True),
        ("                     pytest \u2192 \u2705 \u7da0\u8272", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(8), False),
        ("  IMPROVE  \u91cd\u69cb\uff0c\u4fdd\u6301\u6e2c\u8a66\u7da0\u8272", ACCENT, Pt(24), True),
        ("                     \u6539\u5584\u547d\u540d\u3001\u62bd\u53d6\u5171\u7528\u908f\u8f2f", GRAY_MID, Pt(16), False),
    ], note="看到紅色不要緊張。先定義期望（測試），再寫滿足期望的 code。")

    # FastAPI /docs
    make_content_slide(prs, "FastAPI /docs \u2014 \u514d\u8cbb\u7684\u4e92\u52d5\u6587\u6a94", [
        ("API \u5beb\u5b8c\uff0c\u6587\u6a94\u5c31\u81ea\u52d5\u6709\u4e86", ACCENT, Pt(24), True),
        ("", WHITE, Pt(8), False),
        ("uvicorn app.main:app --reload", WHITE, Pt(20), False),
        ("\u6253\u958b http://localhost:8000/docs", ACCENT2, Pt(20), False),
        ("", WHITE, Pt(12), False),
        ("\u2022  Swagger UI \u81ea\u52d5\u7522\u751f", WHITE, Pt(20), False),
        ("\u2022  \u76f4\u63a5\u5728\u700f\u89bd\u5668\u4e2d\u6e2c\u8a66\u6bcf\u500b API", WHITE, Pt(20), False),
        ("\u2022  \u4e0d\u9700\u8981 Postman \u6216\u5176\u4ed6\u5de5\u5177", WHITE, Pt(20), False),
        ("", WHITE, Pt(12), False),
        ("\u2192  \u9078 FastAPI \u7684\u984d\u5916\u597d\u8655", GRAY_MID, Pt(18), False),
    ], note="Aha Moment：讓學員打開 /docs 頁面，點 Try it out")

    # /review-code
    make_code_slide(prs, "/review-code \u2014 AI \u5be9\u67e5\u7a0b\u5f0f\u78bc", [
        ("\U0001f4cb \u7a0b\u5f0f\u78bc\u5be9\u67e5\u5831\u544a          code-quality-specialist (Sonnet)", ACCENT),
        ("", GRAY_DIM),
        ("\u2705 PASS  \u547d\u540d\u898f\u7bc4", ACCENT2),
        ("\u2705 PASS  \u51fd\u5f0f\u9577\u5ea6 < 50 \u884c", ACCENT2),
        ("\u2705 PASS  \u7121\u786c\u7de8\u78bc\u503c", ACCENT2),
        ("\u26a0\ufe0f  WARN  \u7f3a\u5c11\u8f38\u5165\u9a57\u8b49 (title \u9577\u5ea6)", ACCENT_WARN),
        ("\u2705 PASS  \u932f\u8aa4\u8655\u7406", ACCENT2),
        ("\u2705 PASS  \u4e0d\u53ef\u8b8a\u6a21\u5f0f", ACCENT2),
        ("", GRAY_DIM),
        ("\u6574\u9ad4: 8.5/10", WHITE),
        ("\u5efa\u8b70: \u52a0\u4e0a title \u7684 max_length \u9a57\u8b49", ACCENT_WARN),
    ], note="AI 寫的 code，再讓另一個 AI 審查。依模板的 9 條規則打分。")

    # 第 2 輪
    make_section_slide(prs, "\u7b2c 2 \u8f2a\uff08\u52a0\u901f\uff09", "\u524d\u7aef UI \u2014 React + Tailwind")

    make_content_slide(prs, "\u7b2c 2 \u8f2a\uff1a\u540c\u6a23\u7684\u6d41\u7a0b\uff0c\u52a0\u5feb\u7bc0\u594f", [
        ("/task-next      \u2192  \u81ea\u52d5\u8df3\u5230\u524d\u7aef\u4efb\u52d9", WHITE, Pt(22), False),
        ("/plan + /tdd   \u2192  \u5b78\u54e1\u5df2\u719f\u6089\uff0c\u5feb\u901f\u5e36\u904e", WHITE, Pt(22), False),
        ("", WHITE, Pt(8), False),
        ("\u91cd\u9ede\u89c0\u5bdf\uff1a", ACCENT, Pt(22), True),
        ("", WHITE, Pt(4), False),
        ("\u2022  UI \u8a2d\u8a08\u898f\u5247\u81ea\u52d5\u5957\u7528", WHITE, Pt(20), False),
        ("     \u2192 \u7cfb\u7d71\u5b57\u578b\u3001\u5713\u89d2 12-16px\u3001\u8f15\u5fae\u9670\u5f71\u3001\u5927\u91cf\u7559\u767d", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(4), False),
        ("\u2022  \u4f86\u81ea .claude/rules/ui-design.md", WHITE, Pt(20), False),
        ("     \u2192 \u4e0d\u9700\u8981\u624b\u52d5\u6307\u5b9a\uff0cAI \u81ea\u52d5\u9075\u5b88", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(8), False),
        ("\u2022  \u700f\u89bd\u5668\u4e2d\u67e5\u770b\u6210\u679c\uff1a\u555f\u52d5 dev server\uff0c\u5be6\u969b\u64cd\u4f5c", ACCENT2, Pt(20), False),
    ])

    # 第 3 輪
    make_section_slide(prs, "\u7b2c 3 \u8f2a\uff08\u5feb\u901f\uff09", "\u4fee Bug \u2014 /build-fix")

    make_content_slide(prs, "/build-fix \u2014 \u51fa\u932f\u4e0d\u53ef\u6015", [
        ("\u26a1 \u523b\u610f\u5f15\u5165\u4e00\u500b bug\uff0c\u8b93\u5efa\u7f6e\u5931\u6557", WHITE, Pt(22), False),
        ("", WHITE, Pt(8), False),
        ("/build-fix \u555f\u52d5 build-error-resolver agent (Haiku)", ACCENT, Pt(22), True),
        ("", WHITE, Pt(8), False),
        ("  1.  \u8b80\u53d6\u932f\u8aa4\u8a0a\u606f", WHITE, Pt(20), False),
        ("  2.  \u5b9a\u4f4d\u554f\u984c\u6a94\u6848", WHITE, Pt(20), False),
        ("  3.  \u7528\u6700\u5c0f\u5dee\u7570\u4fee\u5fa9", WHITE, Pt(20), False),
        ("  4.  \u9a57\u8b49\u4fee\u5fa9\u6210\u529f", WHITE, Pt(20), False),
        ("", WHITE, Pt(12), False),
        ("\u8b8a\u66f4: 1 file, 1 line    \u6e2c\u8a66: \u2705 all passed", ACCENT2, Pt(20), False),
        ("", WHITE, Pt(8), False),
        ("\u2192  \u7528\u6700\u8f15\u91cf\u7684 Haiku \u6a21\u578b\uff0c\u5feb\u901f\u4fee\u5fa9\uff0c\u7701\u6210\u672c", GRAY_MID, Pt(18), False),
    ], note="教學效果：刻意犯錯比完美示範更有記憶點")

    # Ch3 小結
    make_content_slide(prs, "Chapter 3 \u5c0f\u7d50 \u2014 3 \u8f2a\u5faa\u74b0\u5b8c\u6210", [
        ("\u7b2c 1 \u8f2a:  \u5f8c\u7aef API    FastAPI + pytest", WHITE, Pt(22), False),
        ("\u7b2c 2 \u8f2a:  \u524d\u7aef UI     React + Vitest", WHITE, Pt(22), False),
        ("\u7b2c 3 \u8f2a:  Bug \u4fee\u5fa9   /build-fix", WHITE, Pt(22), False),
        ("", WHITE, Pt(12), False),
        ("\u6838\u5fc3\u6d41\u7a0b\uff1a", ACCENT, Pt(20), True),
        ("/task-next \u2192 /plan \u2192 /tdd \u2192 /review-code \u2192 commit", ACCENT, Pt(22), False),
        ("", WHITE, Pt(12), False),
        ("Agent \u5354\u4f5c\uff1a", ACCENT2, Pt(20), True),
        ("  planner (Opus) \u2192 tdd-guide (Sonnet) \u2192 code-quality (Sonnet) \u2192 build-resolver (Haiku)", GRAY_LIGHT, Pt(16), False),
    ])

    # ================================================================
    # CHAPTER 4 — 品質驗證
    # ================================================================
    make_section_slide(prs, "Chapter 4", "品質驗證")

    # /verify
    make_code_slide(prs, "/verify \u2014 \u4e09\u9053\u54c1\u8cea\u9580", [
        ("\U0001f50d \u5168\u9762\u9a57\u8b49", ACCENT),
        ("", GRAY_DIM),
        ("Gate 1: \u578b\u5225\u6aa2\u67e5", WHITE),
        ("  Python: mypy --strict     \u2705", ACCENT2),
        ("  TS:     tsc --noEmit      \u2705", ACCENT2),
        ("", GRAY_DIM),
        ("Gate 2: \u6e2c\u8a66", WHITE),
        ("  pytest     \u2192 12/12 passed  \u2705", ACCENT2),
        ("  vitest     \u2192  8/8  passed  \u2705", ACCENT2),
        ("  coverage   \u2192  87%          \u2705 (>80%)", ACCENT2),
        ("", GRAY_DIM),
        ("Gate 3: Lint", WHITE),
        ("  ruff       \u2192  0 errors     \u2705", ACCENT2),
        ("  eslint     \u2192  0 errors     \u2705", ACCENT2),
        ("", GRAY_DIM),
        ("\u7d50\u679c: \u2705 \u5168\u90e8\u901a\u904e", ACCENT2),
    ])

    # /e2e
    make_content_slide(prs, "/e2e \u2014 Playwright \u81ea\u52d5\u64cd\u4f5c\u700f\u89bd\u5668", [
        ("e2e-validation-specialist agent \u6703\uff1a", GRAY_MID, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("  1.  \u555f\u52d5 Playwright \u700f\u89bd\u5668", WHITE, Pt(20), False),
        ("  2.  \u81ea\u52d5\u57f7\u884c\u4f7f\u7528\u8005\u6d41\u7a0b\uff1a", WHITE, Pt(20), False),
        ("", WHITE, Pt(4), False),
        ("        \u6253\u958b\u9996\u9801 \u2192 \u770b\u5230\u7a7a\u7684 Todo \u5217\u8868", GRAY_LIGHT, Pt(18), False),
        ("        \u65b0\u589e \"\u8cb7\u725b\u5976\" \u2192 \u5217\u8868\u51fa\u73fe\u8a72\u9805\u76ee", GRAY_LIGHT, Pt(18), False),
        ("        \u6a19\u8a18\u5b8c\u6210 \u2192 \u51fa\u73fe\u522a\u9664\u7dda", GRAY_LIGHT, Pt(18), False),
        ("        \u5207\u63db\u5206\u985e \u2192 \u7be9\u9078\u6b63\u78ba", GRAY_LIGHT, Pt(18), False),
        ("        \u522a\u9664 \u2192 \u56de\u5230\u7a7a\u5217\u8868", GRAY_LIGHT, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("  3.  \u7522\u51fa\u622a\u5716 + \u6e2c\u8a66\u5831\u544a", ACCENT2, Pt(20), False),
    ], note="視覺衝擊最強的 demo — 瀏覽器自己在動")

    # /time-log
    make_content_slide(prs, "/time-log \u2014 \u81ea\u52d5\u6642\u9593\u8ffd\u8e64", [
        ("\u4f60\u6c92\u6709\u624b\u52d5\u8a18\u904e\u4efb\u4f55\u6642\u9593", ACCENT, Pt(24), True),
        ("", WHITE, Pt(12), False),
        ("  2.1  \u57fa\u790e\u67b6\u69cb              12 min", WHITE, Pt(22), False),
        ("  2.2  CRUD API                18 min", WHITE, Pt(22), False),
        ("  3.1  \u524d\u7aef\u5143\u4ef6              15 min", WHITE, Pt(22), False),
        ("  3.2  CRUD \u4e92\u52d5              12 min", WHITE, Pt(22), False),
        ("  4.1  Bug \u4fee\u5fa9                3 min", WHITE, Pt(22), False),
        ("  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", GRAY_DIM, Pt(22), False),
        ("  \u5408\u8a08                          60 min", ACCENT2, Pt(22), True),
        ("", WHITE, Pt(12), False),
        ("\u2192  \u5168\u90e8\u7531 Hook \u81ea\u52d5\u8ffd\u8e64\uff0c\u5c0d\u56de\u9867\u548c\u672a\u4f86\u4f30\u7b97\u975e\u5e38\u6709\u7528", GRAY_MID, Pt(18), False),
    ], note="用數據收尾：我們花了 X 分鐘完成一個全端應用")

    # ================================================================
    # CHAPTER 5 — 進階 + 收尾
    # ================================================================
    make_section_slide(prs, "Chapter 5", "進階技巧與收尾")

    make_content_slide(prs, "更多可以探索的功能", [
        ("/save-session      \u4fdd\u5b58\u6703\u8a71\u72c0\u614b\uff0c\u4e0b\u6b21\u63a5\u7e8c", WHITE, Pt(22), False),
        ("                             \u2192 \u6bcf\u6b21\u4e0b\u73ed\u524d\u57f7\u884c\uff0c\u9694\u5929\u7121\u7e2b\u63a5\u7e8c", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(6), False),
        ("/hub-delegate       \u624b\u52d5\u9ede\u540d Agent \u57f7\u884c\u4efb\u52d9", WHITE, Pt(22), False),
        ("                             \u2192 \u4f8b\u5982\u76f4\u63a5\u8acb security-auditor \u6383\u63cf\u5b89\u5168", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(6), False),
        ("/learn                   \u8403\u53d6\u53ef\u8907\u7528\u6a21\u5f0f", WHITE, Pt(22), False),
        ("                             \u2192 \u6559 AI \u5b78\u6703\u4f60\u7684\u958b\u767c\u504f\u597d", GRAY_MID, Pt(16), False),
        ("", WHITE, Pt(6), False),
        ("94+ \u5099\u7528\u6280\u80fd\u5305  custom-rule&skill/skills/", WHITE, Pt(22), False),
        ("                             \u2192 \u8907\u88fd\u904e\u4f86\u5c31\u80fd\u7528\uff0c\u6db5\u84cb\u5404\u7a2e\u8a9e\u8a00\u6846\u67b6", GRAY_MID, Pt(16), False),
    ])

    # 流程回顧
    make_content_slide(prs, "完整流程回顧", [
        ("Ch1   \u8a8d\u8b58\u6a21\u677f         5 \u5927\u6a5f\u5236\u300113 \u500b Agent", GRAY_MID, Pt(20), False),
        ("Ch2   /task-init           \u5c08\u6848\u521d\u59cb\u5316 \u2192 CLAUDE.md + WBS", GRAY_MID, Pt(20), False),
        ("Ch3   \u958b\u767c\u5faa\u74b0 x3       /task-next \u2192 /plan \u2192 /tdd \u2192 /review-code", WHITE, Pt(20), False),
        ("Ch4   \u54c1\u8cea\u9a57\u8b49           /verify \u2192 /e2e \u2192 /time-log", WHITE, Pt(20), False),
        ("Ch5   \u9032\u968e\u6280\u5de7           /save-session\u3001/learn\u3001\u64f4\u5145\u6280\u80fd", GRAY_MID, Pt(20), False),
        ("", WHITE, Pt(16), False),
        ("\u2192  ~60 min \u5b8c\u6210\u4e00\u500b\u6709\u6e2c\u8a66\u3001\u6709\u6587\u6a94\u3001\u6709\u54c1\u8cea\u9580\u7684\u5168\u7aef\u61c9\u7528", ACCENT2, Pt(22), True),
    ])

    # 課後作業
    make_content_slide(prs, "課後作業 — 選一個挑戰", [
        ("\u521d\u7d1a", ACCENT2, Pt(24), True),
        ("  \u7528\u6a21\u677f\u5efa\u7acb\u4e00\u500b\u300c\u8a18\u5e33\u672c\u300dApp\uff0c\u8dd1\u5b8c /task-init \u2192 /task-next \u2192 /tdd \u5faa\u74b0", WHITE, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("\u4e2d\u7d1a", ACCENT, Pt(24), True),
        ("  \u70ba TaskFlow \u52a0\u4e0a\u300c\u5230\u671f\u65e5\u300d\u6b04\u4f4d\uff0c\u7df4\u7fd2 SQLModel migration + \u524d\u7aef\u66f4\u65b0", WHITE, Pt(18), False),
        ("", WHITE, Pt(8), False),
        ("\u9032\u968e", ACCENT_WARN, Pt(24), True),
        ("  \u70ba TaskFlow \u52a0\u4e0a JWT \u8a8d\u8b49\uff08FastAPI Security\uff09\uff0c\u9ad4\u9a57 security-auditor agent", WHITE, Pt(18), False),
    ])

    # 探索路徑
    make_content_slide(prs, "推薦探索路徑", [
        ("\u60f3\u6df1\u5165 Agent \u6a5f\u5236     \u2192  .claude/agents/ (13 \u500b\u5b9a\u7fa9\u6a94)", WHITE, Pt(20), False),
        ("\u60f3\u5ba2\u88fd\u5316\u898f\u5247        \u2192  .claude/rules/ (\u4fee\u6539\u6216\u65b0\u589e)", WHITE, Pt(20), False),
        ("\u60f3\u64f4\u5145\u6280\u80fd            \u2192  custom-rule&skill/skills/ (94+)", WHITE, Pt(20), False),
        ("\u60f3\u505a\u5b8c\u6574\u5c08\u6848\u6587\u4ef6  \u2192  /project-docs + 16 \u500b\u6a21\u677f", WHITE, Pt(20), False),
        ("\u60f3\u5efa MCP Server     \u2192  mcp-builder skill", WHITE, Pt(20), False),
        ("\u60f3\u77ad\u89e3\u904b\u4f5c\u539f\u7406     \u2192  .claude/guides/MECHANISMS.md", WHITE, Pt(20), False),
        ("", WHITE, Pt(16), False),
        ("\u901f\u67e5\u5361\uff1a workshop/docs/03_command_cheatsheet.md", ACCENT, Pt(18), False),
    ])

    # Q&A
    s = make_title_slide(prs, "Q & A", "問任何關於模板、工作流、Claude Code 的問題")

    # 感謝
    make_title_slide(prs, "Thank You", "速查卡在 workshop/docs/ \u2014 有問題隨時在內部頻道問")

    return prs


if __name__ == "__main__":
    out_dir = Path(__file__).parent
    prs = build_presentation()
    out_path = out_dir / "VibeCoding_Workshop.pptx"
    prs.save(str(out_path))
    print(f"Generated: {out_path}")
    print(f"Slides: {len(prs.slides)}")
